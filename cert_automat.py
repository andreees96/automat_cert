import pyodbc
#
from pptx import Presentation
#
import win32com.client
win32com.client.gencache.EnsureDispatch('PowerPoint.Application')


connection_string = "Driver={SQL Server};" \
                    "Server=192.xxx.x.xx;" \
                    "Database=GPS;" \
                    "UID=USER;" \
                    "PWD=PASS"

def gps_certificate(patent):
    try:
        ###--LECTURA DE PPT--###
        table_texts = []
        ppt_path = r'C:\Users\Escritorio\certificado.pptx'
        pdf_path = r'C:\Users\Escritorio\certificado_{}.pdf'.format(patent)

        ppt = Presentation(ppt_path)
        slide = ppt.slides[0]
        table = None
        
        
        ##--BUSCAR ELEMENTO TABLA--##
        for element in slide.shapes:
            if element.has_table:
                table = element.table
                break

        for row in table.rows:
            for cell in row.cells:
                table_texts.append(cell.text)

        ###--CONSULTA A LA BD--###
        query = """
                SELECT TOP(1) 
                    CASE
                        WHEN RUT = '99999999-9' THEN ' AUTOMAX - ARRENDADORA DE VEHICULOS S.A.'
                        ELSE '-'
                    END AS RAZON,
                    MOVIL_RUT RUT,
                    REPLACE(LTRIM(RIGHT(CHASIS, LEN(CHASIS) - CHARINDEX('|', CHASIS) - 1)),'ACELEROMETRO','ACELEROMETRO EN 3 EJES') COMPLEMENTOS, 
                    CODIGO PATENTE
                FROM MOVILES WITH(NOLOCK)
                WHERE CODIGO = '{}'""".format(patent)
                
        conn = pyodbc.connect(connection_string)
        cursor = conn.cursor()
        cursor.execute(query)
        row = cursor.fetchone()

        if row:
            table_texts[1] = row.RAZON
            table_texts[3] = row.RUT
            table_texts[5] = row.COMPLEMENTOS
            table_texts[7] = row.PATENTE
        
        print(row)
        conn.close()
        
        ###---ACTUALIZAR DATOS---###
        index_array = 0
        for row in table.rows:
            for cell in row.cells:
                cell.text = table_texts[index_array]
                index_array += 1
                
        ppt.save(ppt_path)
        
        ###---GUARDAR ARCHIVO---###
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(ppt_path, WithWindow=False)
        presentation.Save()
        presentation.ExportAsFixedFormat(pdf_path, 2,PrintRange=None)
        presentation.Close()
        powerpoint.Quit()
        

    except Exception as e:
        print(f"Error: {e}")

# Prueba
gps_certificate('SFXY-93')


